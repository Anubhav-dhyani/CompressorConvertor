from flask import Flask, request, send_file
from werkzeug.utils import secure_filename
from pdf2image import convert_from_path
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from flask_cors import CORS
import os

app = Flask(__name__)
CORS(app)  # <- ADD THIS LINE TO ALLOW CROSS-ORIGIN REQUESTS

UPLOAD_DIR = './uploads'
OUTPUT_DIR = './outputs'

os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)

@app.route('/api/files/convert/pdf-to-doc', methods=['POST'])
def convert_pdf_to_doc():
    file = request.files['file']
    filename = secure_filename(file.filename)
    path = os.path.join(UPLOAD_DIR, filename)
    file.save(path)

    images = convert_from_path(path)
    doc = Document()

    for img in images:
        img_path = os.path.join(OUTPUT_DIR, 'page.jpg')
        img.save(img_path, 'JPEG')
        doc.add_picture(img_path, width=Inches(6))
        doc.add_page_break()

    docx_path = os.path.join(OUTPUT_DIR, filename.replace('.pdf', '.docx'))
    doc.save(docx_path)

    return send_file(docx_path, as_attachment=True)

@app.route('/api/files/convert/pdf-to-ppt', methods=['POST'])
def convert_pdf_to_ppt():
    file = request.files['file']
    filename = secure_filename(file.filename)
    path = os.path.join(UPLOAD_DIR, filename)
    file.save(path)

    images = convert_from_path(path)
    ppt = Presentation()
    blank_slide_layout = ppt.slide_layouts[6]

    for img in images:
        slide = ppt.slides.add_slide(blank_slide_layout)
        img_path = os.path.join(OUTPUT_DIR, 'slide.jpg')
        img.save(img_path, 'JPEG')
        slide.shapes.add_picture(img_path, 0, 0, width=ppt.slide_width, height=ppt.slide_height)

    ppt_path = os.path.join(OUTPUT_DIR, filename.replace('.pdf', '.pptx'))
    ppt.save(ppt_path)

    return send_file(ppt_path, as_attachment=True)


if __name__ == '__main__':
    port = int(os.environ.get("PORT", 5001))
    app.run(host="0.0.0.0", port=port)
